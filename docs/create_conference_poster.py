from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Inches, Pt
import pandas as pd


EMU_PER_INCH = 914400

ACCENT_PURPLE = RGBColor(82, 38, 132)
ACCENT_PURPLE_LIGHT = RGBColor(139, 96, 191)
ACCENT_GOLD = RGBColor(226, 178, 64)
BACKGROUND = RGBColor(249, 246, 253)
SECTION_BG = RGBColor(244, 239, 251)
HIGHLIGHT_BG = RGBColor(254, 247, 226)
TEXT_DARK = RGBColor(33, 30, 47)
TEXT_MUTED = RGBColor(84, 80, 99)
TEAM_LABELS = {
    "UTA": "Utah Jazz",
    "NYK": "New York Knicks",
}


def load_rrs_metrics() -> pd.DataFrame:
    """Load roster resilience metrics with shock deltas for each team-season."""
    rrs = pd.read_csv("vis_clean_data/rrs_from_pos_lineups.csv")
    for col, new_col in [
        ("drop_star", "shock_star"),
        ("drop_role", "shock_role"),
        ("drop_connector", "shock_connector"),
    ]:
        rrs[new_col] = rrs["nr_hat_intact"] - rrs[col]
    rrs_sorted = rrs.sort_values(["season", "team", "RRS"], ascending=[True, True, False])
    return rrs_sorted.drop_duplicates(subset=["season", "team"], keep="first")


def get_team_snapshot(rrs: pd.DataFrame, team: str, season: str | None = None) -> dict:
    """Return a formatted dictionary of snapshot metrics for a team-season."""
    subset = rrs[rrs["team"] == team].copy()
    if subset.empty:
        raise ValueError(f"No roster resilience data for team code '{team}'.")
    if season:
        subset = subset[subset["season"] == season]
        if subset.empty:
            raise ValueError(f"No data for team '{team}' in season '{season}'.")
    snapshot = subset.sort_values("season").iloc[-1]
    return {
        "season": snapshot["season"],
        "team": team,
        "nr_hat_intact": snapshot["nr_hat_intact"],
        "drop_star": snapshot["drop_star"],
        "drop_connector": snapshot["drop_connector"],
        "shock_star": snapshot["shock_star"],
        "shock_role": snapshot["shock_role"],
        "shock_connector": snapshot["shock_connector"],
        "RRS": snapshot["RRS"],
    }


def signed(value: float) -> str:
    return f"{value:+.1f}"


def build_snapshot_items(*snapshots: dict) -> list[dict]:
    items = [
        {
            "text": "Focus teams highlight mesh resilience versus star-centric risk before simulations.",
            "size": 26,
        }
    ]
    for snapshot in snapshots:
        team_name = TEAM_LABELS.get(snapshot["team"], snapshot["team"])
        base = snapshot["nr_hat_intact"]
        star_after = snapshot["drop_star"]
        star_drop = snapshot["shock_star"]
        connector_after = snapshot["drop_connector"]
        connector_drop = snapshot["shock_connector"]
        rrs_value = snapshot["RRS"]
        items.append(
            {
                "text": f"{team_name} {snapshot['season']} (Fig. 1): baseline {signed(base)} NR; star drop {star_drop:.1f}; connector drop {connector_drop:.1f}; RRS {rrs_value:.2f}.",
                "size": 26,
                "bold": True,
                "color": ACCENT_PURPLE,
                "bullet": True,
            }
        )
    return items


def set_no_bullet(paragraph) -> None:
    """Remove bullet styling from a pptx paragraph."""
    p_pr = paragraph._element.get_or_add_pPr()
    for child in list(p_pr):
        if "bu" in child.tag:
            p_pr.remove(child)
    p_pr.append(OxmlElement("a:buNone"))


def add_header(slide, width) -> int:
    header_height = Inches(3.2)
    header = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        left=0,
        top=0,
        width=width,
        height=header_height,
    )
    header.fill.solid()
    header.fill.fore_color.rgb = ACCENT_PURPLE
    header.line.fill.background()

    text_frame = header.text_frame
    text_frame.clear()

    title = text_frame.paragraphs[0]
    title.text = "Roster Geometry & Resilience Across NBA Payroll Networks"
    title.font.size = Pt(88)
    title.font.bold = True
    title.font.color.rgb = RGBColor(255, 255, 255)
    title.alignment = PP_ALIGN.CENTER
    set_no_bullet(title)

    subtitle = text_frame.add_paragraph()
    subtitle.text = (
        "Quantifying how payroll structure, network archetypes, and simulated shocks "
        "shape win equity - Carnegie Mellon Sports Analytics Conference 2025"
    )
    subtitle.font.size = Pt(38)
    subtitle.font.color.rgb = RGBColor(228, 216, 249)
    subtitle.alignment = PP_ALIGN.CENTER
    subtitle.line_spacing = 1.1
    set_no_bullet(subtitle)

    authors = text_frame.add_paragraph()
    authors.text = "Luke Blommesteyn - Lucian Lavric - Yuvraj Sharma"
    authors.font.size = Pt(34)
    authors.font.color.rgb = RGBColor(235, 227, 250)
    authors.alignment = PP_ALIGN.CENTER
    set_no_bullet(authors)

    return header_height + Inches(0.25)


def add_footer(slide, width, height) -> int:
    footer_height = Inches(1.1)
    footer = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        left=0,
        top=height - footer_height,
        width=width,
        height=footer_height,
    )
    footer.fill.solid()
    footer.fill.fore_color.rgb = ACCENT_PURPLE
    footer.line.fill.background()

    text_frame = footer.text_frame
    text_frame.clear()

    note = text_frame.paragraphs[0]
    note.text = "Roster Geometry Project - github.com/16476/csmac - All figures reproducible from open-source pipeline"
    note.font.size = Pt(28)
    note.font.color.rgb = RGBColor(236, 230, 249)
    note.alignment = PP_ALIGN.CENTER
    set_no_bullet(note)

    return footer_height


def add_section_block(
    slide,
    *,
    left,
    top,
    width,
    height,
    heading: str,
    body_items,
    fill_color=SECTION_BG,
    line_color=ACCENT_PURPLE,
    heading_color=ACCENT_PURPLE,
    heading_size=44,
) -> None:
    block = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        left=left,
        top=top,
        width=width,
        height=height,
    )
    block.fill.solid()
    block.fill.fore_color.rgb = fill_color
    block.line.width = Pt(3)
    block.line.color.rgb = line_color
    block.shadow.inherit = False

    text_frame = block.text_frame
    text_frame.clear()
    text_frame.word_wrap = True
    text_frame.margin_left = Pt(16)
    text_frame.margin_right = Pt(16)
    text_frame.margin_top = Pt(20)
    text_frame.margin_bottom = Pt(12)

    heading_para = text_frame.paragraphs[0]
    heading_para.text = heading
    heading_para.font.size = Pt(heading_size)
    heading_para.font.bold = True
    heading_para.font.color.rgb = heading_color
    heading_para.alignment = PP_ALIGN.LEFT
    heading_para.line_spacing = 1
    set_no_bullet(heading_para)

    for item in body_items:
        para = text_frame.add_paragraph()
        para.text = item["text"]
        para.font.size = Pt(item.get("size", 30))
        para.font.bold = item.get("bold", False)
        para.font.italic = item.get("italic", False)
        para.font.color.rgb = item.get("color", TEXT_DARK)
        para.alignment = item.get("alignment", PP_ALIGN.LEFT)
        para.line_spacing = item.get("line_spacing", 1.25)
        if item.get("bullet", False):
            para.level = item.get("level", 0)
        else:
            set_no_bullet(para)


def add_image_with_caption(
    slide,
    image_path: Path,
    *,
    left,
    top,
    width,
    caption: str,
    caption_width=None,
    caption_alignment=PP_ALIGN.CENTER,
) -> int:
    if not image_path.exists():
        return 0

    picture = slide.shapes.add_picture(str(image_path), left=left, top=top, width=width)

    caption_width = caption_width or width
    caption_top = top + picture.height + Inches(0.12)
    caption_height = Inches(0.85)

    caption_box = slide.shapes.add_textbox(
        left=left + (width - caption_width) // 2,
        top=caption_top,
        width=caption_width,
        height=caption_height,
    )
    caption_tf = caption_box.text_frame
    caption_tf.clear()
    paragraph = caption_tf.paragraphs[0]
    paragraph.text = caption
    paragraph.font.size = Pt(24)
    paragraph.font.color.rgb = TEXT_MUTED
    paragraph.font.italic = True
    paragraph.alignment = caption_alignment
    paragraph.line_spacing = 1.1
    set_no_bullet(paragraph)

    return picture.height + caption_height + Inches(0.25)


def add_image_pair_with_caption(
    slide,
    image_paths,
    *,
    left,
    top,
    width,
    gap,
    caption,
    caption_alignment=PP_ALIGN.CENTER,
):
    image_width = (width - gap) // 2
    max_height = 0

    for idx, image_path in enumerate(image_paths):
        path = Path(image_path)
        if not path.exists():
            continue
        pic_left = left + idx * (image_width + gap)
        picture = slide.shapes.add_picture(str(path), left=pic_left, top=top, width=image_width)
        if picture.height > max_height:
            max_height = picture.height

    if max_height == 0:
        return 0

    caption_top = top + max_height + Inches(0.12)
    caption_box = slide.shapes.add_textbox(
        left=left,
        top=caption_top,
        width=width,
        height=Inches(0.9),
    )
    caption_tf = caption_box.text_frame
    caption_tf.clear()
    para = caption_tf.paragraphs[0]
    para.text = caption
    para.font.size = Pt(24)
    para.font.color.rgb = TEXT_MUTED
    para.font.italic = True
    para.alignment = caption_alignment
    para.line_spacing = 1.1
    set_no_bullet(para)

    return max_height + Inches(0.9) + Inches(0.25)


def build_poster(output_path: Path) -> None:
    prs = Presentation()
    prs.slide_width = Inches(48)
    prs.slide_height = Inches(36)

    blank_slide = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide)

    background = slide.background
    background.fill.solid()
    background.fill.fore_color.rgb = BACKGROUND

    top_offset = add_header(slide, prs.slide_width)
    footer_height = add_footer(slide, prs.slide_width, prs.slide_height)

    margin = Inches(0.9)
    gutter = Inches(0.75)
    column_width = (prs.slide_width - (2 * margin) - (2 * gutter)) // 3

    content_top = top_offset + Inches(0.2)
    content_bottom_margin = Inches(0.25)
    column_height = prs.slide_height - content_top - footer_height - content_bottom_margin
    column_gap = Inches(0.28)
    column_bottom = content_top + column_height

    poster_dir = Path(__file__).resolve().parent.parent / "poster_figures"
    rrs_metrics = load_rrs_metrics()
    utah_snapshot = get_team_snapshot(rrs_metrics, "UTA", "2024-2025")
    nyk_snapshot = get_team_snapshot(rrs_metrics, "NYK", "2024-2025")

    # Column 1: Narrative framing
    col1_left = margin
    col1_top = content_top

    intro_items = [
        {
            "text": "We treat roster construction as a network design problem: payroll dollars buy links "
                    "between players, not isolated talent.",
            "size": 28,
        },
        {
            "text": "Fig. 1 compares a Utah \"mesh\" to a New York star core, motivating why we model "
                    "lineup interactions before running shock simulations.",
            "size": 28,
        },
        {
            "text": "Dataset spans eight NBA seasons (240k shared-possession windows) aligned with salaries, "
                    "injuries, and playoff advancement.",
            "size": 28,
        },
    ]
    intro_height = Inches(5.9)
    add_section_block(
        slide,
        left=col1_left,
        top=col1_top,
        width=column_width,
        height=intro_height,
        heading="INTRODUCTION",
        body_items=intro_items,
    )
    col1_top += intro_height + column_gap

    definition_items = [
        {
            "text": "Roster Resilience Score (RRS): average net rating loss across star, role, connector shock templates (lower = stronger).",
            "size": 27,
        },
        {
            "text": "Delta W_s: expected win probability drop for scenario s, calibrated to real injury/usage frequencies.",
            "size": 27,
        },
    ]
    definition_height = Inches(2.9)
    add_section_block(
        slide,
        left=col1_left,
        top=col1_top,
        width=column_width,
        height=definition_height,
        heading="CORE METRICS",
        body_items=definition_items,
        fill_color=HIGHLIGHT_BG,
        line_color=ACCENT_GOLD,
        heading_color=RGBColor(132, 86, 10),
    )
    col1_top += definition_height + column_gap

    network_fig_height = add_image_pair_with_caption(
        slide,
        [
            Path("vis_clean_data") / "fig_network_example1.png",
            Path("vis_clean_data") / "fig_network_example2.png",
        ],
        left=col1_left,
        top=col1_top,
        width=column_width,
        gap=Inches(0.35),
        caption="Fig. 1 Payroll interaction networks: Utah Jazz mesh (left) spreads salary across connectors; New York Knicks core (right) leans on star hubs.",
        caption_alignment=PP_ALIGN.LEFT,
    )
    if network_fig_height:
        col1_top += network_fig_height + column_gap
    else:
        col1_top += column_gap

    snapshot_items = build_snapshot_items(utah_snapshot, nyk_snapshot)
    snapshot_height = Inches(3.9)
    add_section_block(
        slide,
        left=col1_left,
        top=col1_top,
        width=column_width,
        height=snapshot_height,
        heading="TEAM SNAPSHOTS - UTAH VS NEW YORK",
        body_items=snapshot_items,
    )
    col1_top += snapshot_height + column_gap

    questions_items = [
        {
            "text": "Q1: Which roster geometries hold win equity after guard or connector shocks (see Fig. 4)?",
            "size": 26,
            "bullet": True,
        },
        {
            "text": "Q2: How do connector roles and salary assortativity interact with RRS when shocks stack?",
            "size": 26,
            "bullet": True,
        },
        {
            "text": "Q3: What cap-feasible rewires close the gap between fragile and resilient teams?",
            "size": 26,
            "bullet": True,
        },
    ]
    questions_height = Inches(3.8)
    add_section_block(
        slide,
        left=col1_left,
        top=col1_top,
        width=column_width,
        height=questions_height,
        heading="RESEARCH QUESTIONS",
        body_items=questions_items,
    )
    col1_top += questions_height + column_gap

    hypothesis_items = [
        {
            "text": "Mesh salary webs keep >=85% simulated win equity after dual absences (Fig. 4).",
            "size": 27,
        },
        {
            "text": "Star-centric cores without connectors shed win probability fastest (Fig. 5).",
            "size": 27,
        },
    ]
    hypothesis_height = Inches(3.2)
    add_section_block(
        slide,
        left=col1_left,
        top=col1_top,
        width=column_width,
        height=hypothesis_height,
        heading="HYPOTHESIS",
        body_items=hypothesis_items,
        fill_color=HIGHLIGHT_BG,
        line_color=ACCENT_GOLD,
        heading_color=RGBColor(132, 86, 10),
    )
    col1_top += hypothesis_height + column_gap

    data_items = [
        {
            "text": "Sources",
            "bold": True,
            "size": 27,
        },
        {
            "text": "NBA play-by-play + rotation shifts for shared-minute edges; Spotrac salaries; injury and transaction logs.",
            "size": 26,
            "bullet": True,
        },
        {
            "text": "Feature focus",
            "bold": True,
            "size": 27,
        },
        {
            "text": "Connector centrality, salary assortativity, usage entropy, archetype tags feed RRS and Delta W_s modelling.",
            "size": 26,
            "bullet": True,
        },
    ]

    remaining_height = column_bottom - col1_top
    data_height = max(remaining_height, Inches(3.0))
    if col1_top + data_height > column_bottom:
        data_height = column_bottom - col1_top

    add_section_block(
        slide,
        left=col1_left,
        top=col1_top,
        width=column_width,
        height=data_height,
        heading="DATA & FEATURES",
        body_items=data_items,
    )

    # Column 2: Methods and diagnostics
    col2_left = margin + column_width + gutter
    col2_top = content_top

    methods_items = [
        {
            "text": "1. Build roster graphs from shared-minute windows; weight edges by co-playing time and salary flow.",
            "size": 26,
            "bullet": True,
        },
        {
            "text": "2. Derive geometry features (connector centrality, assortativity, usage entropy) feeding RRS and Delta W_s.",
            "size": 26,
            "bullet": True,
        },
        {
            "text": "3. Fig. 2 pipeline: gradient boosted archetype classifier + calibrated logistic regression for win-loss decay.",
            "size": 26,
            "bullet": True,
        },
        {
            "text": "4. Generate 5k cap-feasible synthetic rosters per team to benchmark attainable resilience levels.",
            "size": 26,
            "bullet": True,
        },
    ]
    methods_height = Inches(6.5)
    add_section_block(
        slide,
        left=col2_left,
        top=col2_top,
        width=column_width,
        height=methods_height,
        heading="METHODS",
        body_items=methods_items,
    )
    col2_top += methods_height + column_gap

    workflow_items = [
        {
            "text": "Sample 1-3 player shocks per archetype, weighted by historical injury frequency.",
            "size": 26,
            "bullet": True,
        },
        {
            "text": "Translate each shock into win equity via RAPM-informed lineup forecasts and playoff sims.",
            "size": 26,
            "bullet": True,
        },
        {
            "text": "Log diagnostics (Fig. 3) to check model lift and bootstrap stability across archetypes.",
            "size": 26,
            "bullet": True,
        },
    ]

    workflow_height = Inches(3.0)
    add_section_block(
        slide,
        left=col2_left,
        top=col2_top,
        width=column_width,
        height=workflow_height,
        heading="SIMULATION WORKFLOW",
        body_items=workflow_items,
        fill_color=HIGHLIGHT_BG,
        line_color=ACCENT_GOLD,
        heading_color=RGBColor(132, 86, 10),
    )
    col2_top += workflow_height + column_gap

    cv_height = add_image_with_caption(
        slide,
        poster_dir / "poster_cv_schematic.png",
        left=col2_left,
        top=col2_top,
        width=column_width,
        caption="Fig. 2 Cross-validated modeling pipeline balances geometry features with outcome calibration.",
    )
    if cv_height:
        col2_top += cv_height + Inches(0.2)

    diagnostics_height = add_image_pair_with_caption(
        slide,
        [
            poster_dir / "poster_headline_scorecard.png",
            poster_dir / "poster_bootstrap_macro_f1.png",
        ],
        left=col2_left,
        top=col2_top,
        width=column_width,
        gap=Inches(0.35),
        caption="Fig. 3 Diagnostic panels: scorecard benchmarks lift while bootstrap Macro-F1 curve shows model stability.",
    )
    if diagnostics_height:
        col2_top += diagnostics_height

    # Column 3: Results and implications
    col3_left = margin + (column_width + gutter) * 2
    col3_top = content_top

    results_items = [
        {
            "text": "Fig. 4: Mesh rosters (top-decile connectors) retain ~92% win equity after dual shocks; star cores lose ~18%.",
            "size": 26,
            "bullet": True,
        },
        {
            "text": "Fig. 5a: Connector usage and secondary creator salary share are the strongest levers in the calibrated model.",
            "size": 26,
            "bullet": True,
        },
        {
            "text": "Fig. 5b: Shock-response map shows balanced payroll meshes shrink variance in playoff advancement odds.",
            "size": 26,
            "bullet": True,
        },
    ]
    results_height = Inches(4.8)
    add_section_block(
        slide,
        left=col3_left,
        top=col3_top,
        width=column_width,
        height=results_height,
        heading="RESULTS & INSIGHTS",
        body_items=results_items,
    )
    col3_top += results_height + column_gap

    case_width = Inches(10.8)
    case_left = col3_left + (column_width - case_width) // 2
    case_height = add_image_with_caption(
        slide,
        poster_dir / "poster_case_study_shocks.png",
        left=case_left,
        top=col3_top,
        width=case_width,
        caption="Fig. 4 Case study: mesh roster retains win equity under connector shocks while star-heavy builds collapse.",
        caption_alignment=PP_ALIGN.LEFT,
    )
    if case_height:
        col3_top += case_height + Inches(0.2)
    else:
        col3_top += Inches(0.2)

    feature_pair_height = add_image_pair_with_caption(
        slide,
        [
            poster_dir / "poster_permutation_importance.png",
            poster_dir / "poster_shock_vs_net_rating.png",
        ],
        left=col3_left,
        top=col3_top,
        width=column_width,
        gap=Inches(0.35),
        caption="Fig. 5 Feature importance ranks connector metrics; shock-response map shows how balanced payroll meshes dampen losses.",
        caption_alignment=PP_ALIGN.LEFT,
    )
    if feature_pair_height:
        col3_top += feature_pair_height + column_gap
    else:
        col3_top += column_gap

    conclusion_items = [
        {
            "text": "Takeaways",
            "bold": True,
            "size": 30,
        },
        {
            "text": "Mesh-style payrolls anchored by connector roles are not just aesthetically balanced-they "
                    "are measurably harder to derail.",
            "size": 27,
        },
        {
            "text": "Shifting 6-8% of salary from top earners into secondary creators recovers nearly half of "
                    "the observed resilience gap.",
            "size": 27,
        },
        {
            "text": "Future Work",
            "bold": True,
            "size": 30,
            "color": ACCENT_PURPLE_LIGHT,
        },
        {
            "text": "Integrate physiological load signals, expand archetype embeddings with player tracking data, "
                    "and prototype decision support for front-office scenario planning.",
            "size": 27,
        },
    ]
    conclusion_top = col3_top
    conclusion_height = column_bottom - conclusion_top
    if conclusion_height <= 0:
        conclusion_height = Inches(5.0)
        conclusion_top = column_bottom - conclusion_height

    add_section_block(
        slide,
        left=col3_left,
        top=conclusion_top,
        width=column_width,
        height=conclusion_height,
        heading="CONCLUSION & NEXT STEPS",
        body_items=conclusion_items,
        fill_color=SECTION_BG,
    )

    prs.save(output_path)


if __name__ == "__main__":
    output_file = Path(__file__).resolve().parent / "conference_poster.pptx"
    build_poster(output_file)
